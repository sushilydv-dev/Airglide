import os
import pyautogui
import json
import shutil
import win32com.client  
from tkinter import filedialog  
from pptx import Presentation

screenWidth,screenHeight=pyautogui.size()

def fetch_and_store_pptx():
    pptxFile= filedialog.askopenfilename(filetypes=[("PowerPoint Files","*.pptx")])

    if pptxFile:
        scriptDir= os.path.dirname(os.path.realpath(__file__))

        slidesDir=os.path.join(scriptDir,"pptfile")
        slidesImagesDir=os.path.join(scriptDir,"slides")

#if the dirctory pptfile and slides exists below code will delete them
        if os.path.exists(slidesDir):
            shutil.rmtree(slidesDir)
        if os.path.exists(slidesImagesDir):
            shutil.rmtree(slidesImagesDir)
#the below code will create a new empty directories named pptfile and slides
        os.mkdir(slidesDir)
        os.mkdir(slidesImagesDir)

#nostore the pathpptx flie to be stored in  directory


        pptxFilename=os.path.basename(pptxFile)
        pptxFilePath=os.path.join(slidesDir,pptxFilename)

        shutil.copy(pptxFile,pptxFilePath)
        print(f"file is stored at {pptxFilePath}")
        pptToImageConverter(pptxFilePath, slidesImagesDir)
        print("Slides converted and stored")

        slide_text=extract_text_from_pptx(pptxFilePath)
        text_output_path=os.path.join(scriptDir,"slide_text_data.mjs")
        export_slide_text(slide_text,text_output_path)
        print("text extracted from the slide")


def extract_text_from_pptx(pptx_path):
    prs=Presentation(pptx_path)
    slide_text=[]
    for i , slide in enumerate(prs.slides):
        content=""
        for shape in slide.shapes:
            if hasattr(shape,"text"):
                content+= shape.text.strip()+"\n"
        slide_text.append({
                "slide_number":i+1,
                "text":content.strip()
        })
    return slide_text


import json

def export_slide_text(slide_text, output_file):
    
    json_formatdata = json.dumps(slide_text, indent=4, ensure_ascii=False)
   
    js_variable = f"export const Slide_text_data = {json_formatdata};"

    with open(output_file, "w", encoding="utf-8") as f:
        f.write(js_variable)


def pptToImageConverter(pptxPath, outputDir):
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    powerpoint.Visible = 1

    presentation = powerpoint.Presentations.Open(pptxPath, WithWindow=False)


    for i,slide in enumerate(presentation.Slides):
        slideImagePath=os.path.join(outputDir,f"{i+1}.png")
        slide.Export(slideImagePath,"PNG",screenWidth,screenHeight)
    presentation.Close()
    print(f"Slides exported to {outputDir}")
    powerpoint.Quit()
